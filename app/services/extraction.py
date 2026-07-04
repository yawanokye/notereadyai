from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Iterator

import fitz
from docx import Document
from docx.document import Document as _Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from fastapi import UploadFile
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


@dataclass(slots=True)
class ExtractionResult:
    text: str
    character_count: int
    warning: str | None = None


def _clean_text(text: str) -> str:
    lines = [line.rstrip() for line in text.replace("\r\n", "\n").split("\n")]
    cleaned: list[str] = []
    blank_seen = False
    for line in lines:
        if line.strip():
            cleaned.append(line)
            blank_seen = False
        elif not blank_seen:
            cleaned.append("")
            blank_seen = True
    return "\n".join(cleaned).strip()


def _extract_pdf(data: bytes) -> str:
    document = fitz.open(stream=data, filetype="pdf")
    pages = []
    for page_number, page in enumerate(document, start=1):
        text = page.get_text("text").strip()
        if text:
            pages.append(f"## Page {page_number}\n\n{text}")
    return "\n\n".join(pages)


def _iter_docx_blocks(document: _Document) -> Iterator[Paragraph | Table]:
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _escape_cell(value: str) -> str:
    return " ".join(value.replace("|", "\\|").split())


def _table_to_markdown(table: Table) -> str:
    rows = [[_escape_cell(cell.text) for cell in row.cells] for row in table.rows]
    rows = [row for row in rows if any(row)]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalised = [row + [""] * (width - len(row)) for row in rows]
    header = normalised[0]
    separator = ["---"] * width
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in normalised[1:])
    return "\n".join(lines)


def _extract_docx(data: bytes) -> str:
    document = Document(BytesIO(data))
    parts: list[str] = []
    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style = (block.style.name or "").lower()
            if style.startswith("heading"):
                try:
                    level = min(4, max(1, int(style.split()[-1])))
                except (ValueError, IndexError):
                    level = 2
                parts.append(f"{'#' * level} {text}")
            elif "list bullet" in style:
                parts.append(f"- {text}")
            elif "list number" in style:
                parts.append(f"1. {text}")
            else:
                parts.append(text)
        else:
            markdown_table = _table_to_markdown(block)
            if markdown_table:
                parts.append(markdown_table)
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    presentation = Presentation(BytesIO(data))
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        items: list[str] = []
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text.strip():
            items.append(f"## {title_shape.text.strip()}")
        else:
            items.append(f"## Slide {index}")
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = [[_escape_cell(cell.text) for cell in row.cells] for row in table.rows]
                if rows:
                    width = max(len(row) for row in rows)
                    rows = [row + [""] * (width - len(row)) for row in rows]
                    items.append("| " + " | ".join(rows[0]) + " |")
                    items.append("| " + " | ".join(["---"] * width) + " |")
                    items.extend("| " + " | ".join(row) + " |" for row in rows[1:])
            elif hasattr(shape, "text") and shape.text.strip():
                items.append(shape.text.strip())
        slides.append("\n\n".join(items))
    return "\n\n".join(slides)


async def extract_upload(upload: UploadFile, max_chars: int) -> ExtractionResult:
    filename = upload.filename or "uploaded-file"
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type. Upload one of: {supported}.")

    data = await upload.read()
    if not data:
        raise ValueError("The uploaded file is empty.")

    if extension == ".pdf":
        text = _extract_pdf(data)
    elif extension == ".docx":
        text = _extract_docx(data)
    elif extension == ".pptx":
        text = _extract_pptx(data)
    else:
        text = data.decode("utf-8", errors="replace")

    text = _clean_text(text)
    if not text:
        raise ValueError(
            "No readable text was extracted. The file may contain scanned pages. "
            "OCR support is not yet enabled."
        )

    original_count = len(text)
    warning = None
    if original_count > max_chars:
        text = text[:max_chars]
        warning = (
            f"The source contained {original_count:,} characters. The first "
            f"{max_chars:,} characters were processed. Split very large source files "
            "or increase MAX_EXTRACTED_CHARS."
        )

    return ExtractionResult(text=text, character_count=len(text), warning=warning)
